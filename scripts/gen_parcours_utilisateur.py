"""
Génère UC28_Parcours_Utilisateur.pptx — deck du parcours utilisateur
basé sur le Script_Demo_UC28_Opus48.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Palette ──────────────────────────────────────────────────────────────────
C_DARK_TEAL   = RGBColor(0x00, 0x3C, 0x50)   # fond header
C_BRAND       = RGBColor(0x00, 0x70, 0xAD)   # bleu BV
C_CYAN        = RGBColor(0x00, 0xB0, 0xF0)   # cyan accent
C_EMERALD     = RGBColor(0x00, 0x7E, 0x5C)   # vert émeraude
C_AMBER       = RGBColor(0xD9, 0x7A, 0x00)   # ambre alerte
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG    = RGBColor(0xEE, 0xF2, 0xF7)   # fond canvas
C_INK         = RGBColor(0x1A, 0x2E, 0x3B)   # texte sombre
C_INK_MUTED   = RGBColor(0x60, 0x7D, 0x8B)   # texte gris
C_SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)   # card white
C_DIVIDER     = RGBColor(0xCF, 0xD8, 0xDC)
C_RED         = RGBColor(0xC0, 0x39, 0x2B)
C_SECTION_BG  = RGBColor(0xE3, 0xF2, 0xFD)   # wash cyan

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT_PATH      = os.path.join(os.path.dirname(__file__), "..", "uc28-inspection", "docs", "UC28_Parcours_Utilisateur.pptx")
SCREENSHOTS   = os.path.join(os.path.dirname(__file__), "..", "uc28-inspection", "docs", "screenshots")


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text_box(slide, l, t, w, h, text, size=Pt(12), bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb


def add_para(tf, text, size=Pt(11), bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return p


def header_bar(slide, title, subtitle=None, icon_text="BV"):
    """Barre d'en-tête dark-teal sur toute la largeur."""
    bar_h = Inches(1.1)
    rect(slide, 0, 0, SLIDE_W, bar_h, fill=C_DARK_TEAL)

    # Logo carré blanc
    logo = rect(slide, Inches(0.3), Inches(0.25), Inches(0.6), Inches(0.6), fill=C_WHITE)
    text_box(slide, Inches(0.3), Inches(0.30), Inches(0.6), Inches(0.5),
             icon_text, size=Pt(13), bold=True, color=C_DARK_TEAL, align=PP_ALIGN.CENTER)

    # Titre
    text_box(slide, Inches(1.1), Inches(0.18), Inches(9), Inches(0.45),
             title, size=Pt(18), bold=True, color=C_WHITE)
    if subtitle:
        text_box(slide, Inches(1.1), Inches(0.62), Inches(9), Inches(0.35),
                 subtitle, size=Pt(10), color=C_CYAN, italic=True)


def page_bg(slide):
    """Fond canvas clair."""
    rect(slide, 0, Inches(1.1), SLIDE_W, SLIDE_H - Inches(1.1), fill=C_LIGHT_BG)


def card(slide, l, t, w, h, fill=C_SURFACE, border=C_DIVIDER):
    rect(slide, l, t, w, h, fill=fill, line_color=border, line_w=Pt(0.75))


def section_label(slide, l, t, w, text, color=C_BRAND):
    """Label de section en majuscules petits."""
    text_box(slide, l, t, w, Inches(0.25), text.upper(),
             size=Pt(8), bold=True, color=color)


def bullet_box(slide, l, t, w, h, items, title=None, title_color=C_BRAND,
               bullet="•", item_size=Pt(11), item_color=C_INK,
               fill=C_SURFACE, border=C_DIVIDER):
    card(slide, l, t, w, h, fill=fill, border=border)
    top = t + Inches(0.15)
    if title:
        text_box(slide, l + Inches(0.15), top, w - Inches(0.3), Inches(0.3),
                 title, size=Pt(10), bold=True, color=title_color)
        top += Inches(0.32)
    txb = slide.shapes.add_textbox(l + Inches(0.15), top, w - Inches(0.3),
                                   h - (top - t) - Inches(0.1))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = item_size
        run.font.color.rgb = item_color
    return txb


def step_badge(slide, l, t, label, color=C_BRAND):
    """Petit badge arrondi simulé avec un rectangle coloré."""
    rect(slide, l, t, Inches(2.0), Inches(0.3), fill=color)
    text_box(slide, l, t, Inches(2.0), Inches(0.3), label,
             size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def acte_badge(slide, l, t, acte, label, color=C_DARK_TEAL):
    rect(slide, l, t, Inches(2.4), Inches(0.38), fill=color)
    text_box(slide, l, t, Inches(2.4), Inches(0.38),
             f"{acte}  ·  {label}",
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE FACTORY
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    """Slide 1 — Couverture"""
    s = blank(prs)

    # Fond pleine page dark-teal
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK_TEAL)

    # Bandeau cyan bas
    rect(s, 0, SLIDE_H - Inches(0.6), SLIDE_W, Inches(0.6), fill=C_BRAND)

    # Logo BV
    rect(s, Inches(0.5), Inches(0.5), Inches(0.85), Inches(0.85), fill=C_WHITE)
    text_box(s, Inches(0.5), Inches(0.55), Inches(0.85), Inches(0.75),
             "BV", size=Pt(16), bold=True, color=C_DARK_TEAL, align=PP_ALIGN.CENTER)

    # Surtitle
    text_box(s, Inches(1.55), Inches(0.55), Inches(9), Inches(0.4),
             "UC 28 · HACKATHON CAPGEMINI × ANTHROPIC 2026",
             size=Pt(10), color=C_CYAN, italic=True)

    # Titre principal
    text_box(s, Inches(0.5), Inches(1.6), Inches(12), Inches(1.0),
             "Inspection Augmentée",
             size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Sous-titre — parcours
    text_box(s, Inches(0.5), Inches(2.65), Inches(12), Inches(0.5),
             "Parcours utilisateur — La journée de Marc Lefèvre",
             size=Pt(20), color=C_CYAN, italic=True)

    # Ligne de séparation
    rect(s, Inches(0.5), Inches(3.3), Inches(5), Inches(0.04), fill=C_CYAN)

    # Accroche
    txb = s.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(10), Inches(1.8))
    tf  = txb.text_frame
    tf.word_wrap = True
    for line, sz, clr in [
        ('"Il est 21 heures.', Pt(15), C_WHITE),
        (' Marc a terminé son audit à 16 heures.', Pt(15), C_WHITE),
        (' Et il est encore devant son écran."', Pt(15), C_WHITE),
        ("", Pt(10), C_WHITE),
        ("Nous, on a inversé ce ratio.", Pt(13), C_CYAN),
    ]:
        p = tf.paragraphs[0] if line == '"Il est 21 heures.' else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size  = sz
        run.font.color.rgb = clr
        run.font.italic = (clr == C_CYAN)

    # Équipe
    text_box(s, Inches(0.5), Inches(5.5), Inches(8), Inches(0.4),
             "Équipe Code Resonance — Philippe Gragnic · Habib Koffi · Véronique Poilane Zhang",
             size=Pt(10), color=C_INK_MUTED)

    # Durée
    text_box(s, SLIDE_W - Inches(3.2), Inches(5.5), Inches(2.8), Inches(0.4),
             "⏱  8 minutes · 3 actes + épilogue",
             size=Pt(10), color=C_INK_MUTED, align=PP_ALIGN.RIGHT)


def slide_douleurs(prs):
    """Slide 2 — Les 5 douleurs"""
    s = blank(prs)
    header_bar(s, "Le problème — Cinq douleurs, une seule journée",
               "Bureau Veritas · Auditeurs ISO 9001 · Marc Lefèvre, 23 ans de terrain")
    page_bg(s)

    TOP = Inches(1.3)
    col_w = Inches(2.4)
    pad = Inches(0.18)

    DOULEURS = [
        ("1", "Expertise rare", "C_BRAND",
         ["Savoir-faire sénior difficile\nà transmettre", "Part avec l'auditeur\nà la retraite"]),
        ("2", "Préparation longue", "C_AMBER",
         ["Hétérogène d'un\nauditeur à l'autre", "Sources multiples,\npas de standard"]),
        ("3", "Historique ignoré", "C_BRAND",
         ["NC passées jamais\nrecroisées sur terrain", "Récidives non\ndétectées"]),
        ("4", "Constats subjectifs", "C_AMBER",
         ["Formulation variable,\npas de référentiel", "Clause ISO citée\nde mémoire"]),
        ("5", "Rédaction chronophage", "C_RED",
         ["4 h de rapport\naprès 8 h d'audit", "Soirées sacrifiées,\nqualité inégale"]),
    ]

    colors = {
        "C_BRAND": C_BRAND, "C_AMBER": C_AMBER, "C_RED": C_RED
    }

    for i, (num, titre, color_key, items) in enumerate(DOULEURS):
        cx = Inches(0.35) + i * (col_w + pad)
        c  = colors[color_key]

        # Card
        card(s, cx, TOP, col_w, Inches(4.2), fill=C_SURFACE, border=C_DIVIDER)

        # Numéro
        rect(s, cx, TOP, col_w, Inches(0.45), fill=c)
        text_box(s, cx, TOP, col_w, Inches(0.45),
                 f"Douleur {num}", size=Pt(11), bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

        # Titre
        text_box(s, cx + Inches(0.12), TOP + Inches(0.55), col_w - Inches(0.24), Inches(0.55),
                 titre, size=Pt(13), bold=True, color=C_INK)

        # Items
        txb = s.shapes.add_textbox(cx + Inches(0.12), TOP + Inches(1.15),
                                   col_w - Inches(0.24), Inches(2.8))
        tf  = txb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.size  = Pt(10)
            run.font.color.rgb = C_INK_MUTED

    # Bas — promesse
    rect(s, Inches(0.35), Inches(5.8), SLIDE_W - Inches(0.7), Inches(0.9), fill=C_BRAND)
    text_box(s, Inches(0.35), Inches(5.9), SLIDE_W - Inches(0.7), Inches(0.7),
             "→  Une seule journée de démo pour traiter ces cinq douleurs — avant, pendant, et après l'audit.",
             size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Slide 3 — La boucle fermée"""
    s = blank(prs)
    header_bar(s, "La boucle fermée — Une plateforme, deux faces",
               "Bureau Veritas achète · RATP & Apave audités · Moteur IA partagé")
    page_bg(s)

    TOP = Inches(1.3)

    # ── Face BV ──────────────────────────────────────────────────
    card(s, Inches(0.4), TOP, Inches(3.8), Inches(4.6), fill=C_SURFACE, border=C_BRAND)
    rect(s, Inches(0.4), TOP, Inches(3.8), Inches(0.5), fill=C_BRAND)
    text_box(s, Inches(0.4), TOP + Inches(0.05), Inches(3.8), Inches(0.4),
             "👤  MARC LEFÈVRE — Auditeur BV", size=Pt(11), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

    bv_items = [
        "Dashboard : 10 missions planifiées",
        "Brief auto-généré (site + historique)",
        "Check-list ISO adaptée au scope",
        "Capture : micro + photo + analyse IA",
        "Rapport instantané sur site",
        "Transmission rapport → portail RATP",
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.65), Inches(3.4), Inches(3.7))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bv_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(11)
        run.font.color.rgb = C_INK

    # ── Centre — moteur IA ────────────────────────────────────────
    cx = Inches(4.5)
    card(s, cx, TOP + Inches(0.6), Inches(4.3), Inches(3.5), fill=C_SECTION_BG, border=C_CYAN)
    text_box(s, cx + Inches(0.1), TOP + Inches(0.75), Inches(4.1), Inches(0.4),
             "MOTEUR IA PARTAGÉ", size=Pt(10), bold=True, color=C_BRAND, align=PP_ALIGN.CENTER)

    ia_items = [
        ("Claude Opus", "via GEP Capgemini"),
        ("RAG ISO 9001:2015", "sentence-transformers BERT"),
        ("Analyse constat", "clause + criticité + actions"),
        ("Pré-analyse documents", "fournisseur avant l'audit"),
        ("Explicabilité", "extrait verbatim de la norme"),
        ("Capitalisation", "feedbacks Confirmer/Corriger"),
    ]
    txb = s.shapes.add_textbox(cx + Inches(0.2), TOP + Inches(1.25), Inches(3.9), Inches(2.7))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(ia_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▸  {k}"
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = f"  {v}"
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK_MUTED

    # Flèches
    text_box(s, Inches(4.15), TOP + Inches(1.8), Inches(0.5), Inches(0.5),
             "→", size=Pt(24), bold=True, color=C_BRAND, align=PP_ALIGN.CENTER)
    text_box(s, Inches(8.7), TOP + Inches(1.8), Inches(0.5), Inches(0.5),
             "→", size=Pt(24), bold=True, color=C_EMERALD, align=PP_ALIGN.CENTER)
    text_box(s, Inches(8.7), TOP + Inches(2.6), Inches(0.5), Inches(0.5),
             "←", size=Pt(24), bold=True, color=C_EMERALD, align=PP_ALIGN.CENTER)

    # ── Face RATP ─────────────────────────────────────────────────
    card(s, Inches(9.1), TOP, Inches(3.8), Inches(4.6), fill=C_SURFACE, border=C_EMERALD)
    rect(s, Inches(9.1), TOP, Inches(3.8), Inches(0.5), fill=C_EMERALD)
    text_box(s, Inches(9.1), TOP + Inches(0.05), Inches(3.8), Inches(0.4),
             "🏢  MEI LIN ZHANG — Qualité RATP", size=Pt(11), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

    ratp_items = [
        "Dépôt de documents (procédures, CR)",
        "Pré-analyse Claude avant l'audit",
        "Alertes ⚠ remontées dans le brief BV",
        "Réception du rapport signé",
        "Badge Bureau Veritas + analyse NC",
        "Téléchargement du Word signé",
    ]
    txb = s.shapes.add_textbox(Inches(9.3), TOP + Inches(0.65), Inches(3.4), Inches(3.7))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(ratp_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(11)
        run.font.color.rgb = C_INK

    # Bas — slogan boucle
    rect(s, Inches(0.4), Inches(6.1), SLIDE_W - Inches(0.8), Inches(0.7), fill=C_DARK_TEAL)
    text_box(s, Inches(0.4), Inches(6.2), SLIDE_W - Inches(0.8), Inches(0.5),
             "RATP dépose  →  BV audite  →  le rapport revient à RATP  ·  La boucle est complète",
             size=Pt(13), bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)


def slide_connexion(prs):
    """Slide 4 — Écran 0 : Connexion"""
    s = blank(prs)
    header_bar(s, "Écran 0 — Connexion", "Deux rôles, une seule application")
    page_bg(s)

    TOP = Inches(1.4)

    # Carte Marc
    card(s, Inches(0.7), TOP, Inches(5.5), Inches(3.8), fill=C_SURFACE, border=C_BRAND)
    rect(s, Inches(0.7), TOP, Inches(5.5), Inches(0.5), fill=C_BRAND)
    text_box(s, Inches(0.7), TOP + Inches(0.05), Inches(5.5), Inches(0.4),
             "👤  Marc Lefèvre — Auditeur BV",
             size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    marc_items = [
        "Profil pré-chargé — aucune saisie",
        "Connexion en un clic",
        "→ ClientList → Dashboard → Brief → Inspection → Rapport",
        "Parcours principal : 8 minutes",
    ]
    txb = s.shapes.add_textbox(Inches(0.9), TOP + Inches(0.65), Inches(5.1), Inches(2.9))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(marc_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(11)
        run.font.color.rgb = C_INK

    # Carte Mei Lin
    card(s, Inches(7.1), TOP, Inches(5.5), Inches(3.8), fill=C_SURFACE, border=C_EMERALD)
    rect(s, Inches(7.1), TOP, Inches(5.5), Inches(0.5), fill=C_EMERALD)
    text_box(s, Inches(7.1), TOP + Inches(0.05), Inches(5.5), Inches(0.4),
             "🏢  Mei Lin Zhang — Qualité RATP",
             size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    mei_items = [
        "Accès direct au portail fournisseur",
        "Dépose procédures, CR, plans qualité",
        "Reçoit le rapport BV en fin de boucle",
        "Rôle : épilogue — fermer la boucle",
    ]
    txb = s.shapes.add_textbox(Inches(7.3), TOP + Inches(0.65), Inches(5.1), Inches(2.9))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(mei_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(11)
        run.font.color.rgb = C_INK

    # Bas — note
    rect(s, Inches(0.7), Inches(5.45), SLIDE_W - Inches(1.4), Inches(0.55), fill=C_SECTION_BG)
    text_box(s, Inches(0.9), Inches(5.52), SLIDE_W - Inches(1.8), Inches(0.45),
             "ℹ  On suit Marc tout au long de la démo. Mei Lin Zhang revient en Épilogue pour refermer la boucle.",
             size=Pt(10), color=C_BRAND, italic=True)


def slide_dashboard(prs):
    """Slide 5 — Écran 1 : Dashboard"""
    s = blank(prs)
    header_bar(s, "Écran 1 — Dashboard", "La journée de Marc d'un seul regard")
    page_bg(s)

    TOP = Inches(1.35)

    # KPIs — ligne
    kpis = [
        ("10", "Audits planifiés", C_BRAND),
        ("2", "TERMINÉS", C_EMERALD),
        ("1", "EN COURS", C_AMBER),
        ("7", "PLANIFIÉS", C_INK_MUTED),
    ]
    for i, (val, label, c) in enumerate(kpis):
        cx = Inches(0.4) + i * Inches(3.2)
        card(s, cx, TOP, Inches(3.0), Inches(0.95), fill=C_SURFACE, border=C_DIVIDER)
        text_box(s, cx + Inches(0.1), TOP + Inches(0.04), Inches(1.1), Inches(0.87),
                 val, size=Pt(26), bold=True, color=c, align=PP_ALIGN.CENTER)
        text_box(s, cx + Inches(1.2), TOP + Inches(0.3), Inches(1.6), Inches(0.55),
                 label, size=Pt(10), color=C_INK_MUTED)

    TOP2 = TOP + Inches(1.1)

    # Col gauche — fonctionnalités
    card(s, Inches(0.4), TOP2, Inches(6.0), Inches(4.0), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP2 + Inches(0.15), Inches(5.6), Inches(0.35),
             "FONCTIONNALITÉS", size=Pt(9), bold=True, color=C_INK_MUTED)

    features = [
        ("Trait rouge temps réel", "Heure courante sur la timeline — où en est Marc ?"),
        ("Vue Liste + Planning", "Bascule entre vue cartes et créneaux horodatés"),
        ("Filtre statuts", "TERMINÉ / PROCHAIN / PLANIFIÉ — filtrage en un clic"),
        ("Carte Leaflet + OSRM", "Itinéraire routier réel, clic → zoom sur le site"),
        ("Temps de trajet", "Voiture / TC / Vélo intégrés dans chaque carte"),
        ("Démarrer l'inspection", "Bouton vert sur la mission PROCHAIN : Sucy-en-Brie"),
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP2 + Inches(0.55), Inches(5.6), Inches(3.3))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"▸  {k}  "
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(11)
        run2.font.color.rgb = C_INK_MUTED

    # Col droite — scénario
    card(s, Inches(6.7), TOP2, Inches(6.2), Inches(4.0), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(6.9), TOP2 + Inches(0.15), Inches(5.8), Inches(0.35),
             "SCÉNARIO", size=Pt(9), bold=True, color=C_INK_MUTED)

    steps = [
        "Clic sur une carte → carte zoomée sur le site",
        "Bascule Planning → voir les créneaux + trajets",
        "Filtre PROCHAIN → identifier Sucy-en-Brie",
        "Clic Démarrer → lancer le Brief",
    ]
    txb = s.shapes.add_textbox(Inches(6.9), TOP2 + Inches(0.55), Inches(5.8), Inches(2.2))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"{i+1}.  {step}"
        run.font.size  = Pt(12)
        run.font.color.rgb = C_INK

    # Highlight mission Sucy
    rect(s, Inches(6.9), TOP2 + Inches(2.9), Inches(5.8), Inches(0.8), fill=C_SECTION_BG, line_color=C_BRAND, line_w=Pt(1))
    text_box(s, Inches(7.05), TOP2 + Inches(3.0), Inches(5.5), Inches(0.65),
             "Mission cible : Atelier Sucy-en-Brie — Maintenance rames MI09 RER A\n§7.1.5 Métrologie · §8.7 Non-conformités · 2 h 30 · PROCHAIN",
             size=Pt(10), color=C_BRAND)


def slide_brief_1(prs):
    """Slide 6 — Acte I : Brief (contexte + historique)"""
    s = blank(prs)
    header_bar(s, "ACTE I — Brief pré-audit", "Avant de quitter son bureau — T 2:20")
    page_bg(s)

    TOP = Inches(1.35)
    acte_badge(s, Inches(0.4), TOP, "ACTE I", "AVANT", C_BRAND)
    TOP += Inches(0.5)

    # Col gauche — contexte site
    card(s, Inches(0.4), TOP, Inches(6.0), Inches(2.2), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(0.12), Inches(5.6), Inches(0.32),
             "CONTEXTE SITE — chargé depuis la base BV", size=Pt(9), bold=True, color=C_INK_MUTED)
    ctx_items = [
        ("Site", "Atelier Sucy-en-Brie · RATP-SUC"),
        ("Périmètre", "Maintenance rames MI09 — RER A"),
        ("Effectif", "218 personnes"),
        ("Responsable qualité", "Mei Lin Zhang"),
        ("Source", "GET /api/sites/RATP-SUC → audit.db"),
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.5), Inches(5.6), Inches(1.6))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(ctx_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{k} : "
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK

    # Col droite — historique + alerte
    card(s, Inches(6.7), TOP, Inches(6.2), Inches(2.2), fill=C_SURFACE, border=C_RED)
    rect(s, Inches(6.7), TOP, Inches(6.2), Inches(0.4), fill=C_RED)
    text_box(s, Inches(6.7), TOP + Inches(0.04), Inches(6.2), Inches(0.33),
             "⚠  HISTORIQUE — NC NON CLÔTURÉE", size=Pt(10), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    hist_items = [
        "Nov. 2024 : NC MAJEURE §7.1.5 — certificats d'étalonnage expirés",
        "NC toujours ouverte — jamais clôturée",
        "Action corrective : étalonnage COFRAC (non réalisé)",
        "→ Section S2 pré-flagguée, priorité haute dans la check-list",
        "Thèmes récurrents : §7.1.5 + §8.7 → badge 🔁 en inspection",
    ]
    txb = s.shapes.add_textbox(Inches(6.9), TOP + Inches(0.5), Inches(5.8), Inches(1.6))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(hist_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size  = Pt(10)
        run.font.color.rgb = C_INK

    # Bas — check-list auto-générée (résumé)
    card(s, Inches(0.4), TOP + Inches(2.35), SLIDE_W - Inches(0.8), Inches(2.6),
         fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(2.5), Inches(12), Inches(0.35),
             "CHECK-LIST AUTO-GÉNÉRÉE — 5 sources combinées", size=Pt(9), bold=True, color=C_INK_MUTED)

    sections = [
        ("S1", "§7.5", "Maîtrise documentaire", C_INK_MUTED, "générique"),
        ("S2", "§7.1.5", "Étalonnage & équipements de mesure", C_RED, "NC ouverte + récurrence 🔁"),
        ("S3", "§8.7", "Gestion des non-conformités", C_AMBER, "scope + récurrence 🔁"),
        ("S4", "§7.2", "Compétences", C_INK_MUTED, "lié à 218 personnes"),
    ]
    for i, (num, clause, titre, c, note) in enumerate(sections):
        cx = Inches(0.6) + i * Inches(3.15)
        rect(s, cx, TOP + Inches(2.9), Inches(2.95), Inches(0.32), fill=c)
        text_box(s, cx, TOP + Inches(2.9), Inches(2.95), Inches(0.32),
                 f"{num}  {clause}  {titre}", size=Pt(9), bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        text_box(s, cx, TOP + Inches(3.27), Inches(2.95), Inches(0.4),
                 note, size=Pt(9), color=c, align=PP_ALIGN.CENTER, italic=True)


def slide_brief_2(prs):
    """Slide 7 — Acte I : Brief (scope + docs RATP + édition check-list)"""
    s = blank(prs)
    header_bar(s, "ACTE I — Brief pré-audit (suite)",
               "Scope, documents fournisseur, check-list éditable — Douleurs 1 & 2")
    page_bg(s)

    TOP = Inches(1.35)

    # Col 1 — Scope de l'audit
    card(s, Inches(0.4), TOP, Inches(4.0), Inches(5.0), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(0.15), Inches(3.6), Inches(0.35),
             "SCOPE DE L'AUDIT", size=Pt(9), bold=True, color=C_INK_MUTED)
    scope_items = [
        "✓  Étalonnage & équipements de mesure (§7.1.5)",
        "✓  Gestion des non-conformités (§8.7)",
        "✓  Maîtrise documentaire (§7.5)",
        "",
        "+ Ajouter un scope →",
        "Élargir le périmètre si besoin",
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.6), Inches(3.6), Inches(3.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(scope_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(10)
        run.font.color.rgb = C_BRAND if item.startswith("+") else (C_INK_MUTED if not item else C_INK)

    # Col 2 — Docs portail RATP
    card(s, Inches(4.65), TOP, Inches(4.0), Inches(5.0), fill=C_SURFACE, border=C_AMBER)
    rect(s, Inches(4.65), TOP, Inches(4.0), Inches(0.4), fill=C_AMBER)
    text_box(s, Inches(4.65), TOP + Inches(0.04), Inches(4.0), Inches(0.33),
             "⚠  DOCUMENTS PORTAIL RATP", size=Pt(10), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    docs = [
        ("✓", "CR_Audit_nov_2024.docx", "14/11/2024 · §7.1.5 signalé"),
        ("✓", "Procedures_etalonnage_v3.pdf", "v3 · 3 équipements périmés"),
        ("✓", "Plan_qualite_2025.pdf", "Budget étalonnage : non alloué"),
    ]
    txb = s.shapes.add_textbox(Inches(4.85), TOP + Inches(0.55), Inches(3.6), Inches(2.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (icon, nom, note) in enumerate(docs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"{icon}  {nom}"
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = C_EMERALD
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = f"     {note}"
        run2.font.size  = Pt(9)
        run2.font.color.rgb = C_INK_MUTED

    text_box(s, Inches(4.85), TOP + Inches(3.3), Inches(3.6), Inches(0.9),
             "⚠  §7.1.5 et §8.7 marqués en alerte\ndans la check-list de Marc",
             size=Pt(10), color=C_AMBER)

    # Col 3 — Check-list éditable
    card(s, Inches(8.9), TOP, Inches(4.0), Inches(5.0), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(9.1), TOP + Inches(0.15), Inches(3.6), Inches(0.35),
             "CHECK-LIST ÉDITABLE", size=Pt(9), bold=True, color=C_INK_MUTED)
    edit_items = [
        ("× sur un item →", "supprimer un point inutile"),
        ("× sur une section →", "exclure une section entière"),
        ("+ Ajouter un point →", "créer un point sur mesure"),
        ("Badge Auditeur", "le point porte son nom — traçable"),
        ("", ""),
        ("→ Ces ajouts entrent", "dans le système :"),
        ("capitalisation", "de l'expertise BV"),
    ]
    txb = s.shapes.add_textbox(Inches(9.1), TOP + Inches(0.6), Inches(3.6), Inches(3.8))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(edit_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{k}  "
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND if k else C_EMERALD
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK_MUTED

    # Douleurs traitées
    rect(s, Inches(0.4), Inches(6.55), SLIDE_W - Inches(0.8), Inches(0.65), fill=C_BRAND)
    text_box(s, Inches(0.6), Inches(6.62), SLIDE_W - Inches(1.2), Inches(0.5),
             "Douleurs traitées :  ① Expertise capitalisée (badge Auditeur)  ·  ② Préparation standardisée (check-list auto-générée)",
             size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_inspection_1(prs):
    """Slide 8 — Acte II : Capture sur le terrain"""
    s = blank(prs)
    header_bar(s, "ACTE II — Inspection sur le terrain",
               "Marc est sur site — Gants aux mains, devant les équipements — T 3:30")
    page_bg(s)

    TOP = Inches(1.35)
    acte_badge(s, Inches(0.4), TOP, "ACTE II", "PENDANT", C_EMERALD)
    TOP += Inches(0.5)

    # Bannière récurrences
    rect(s, Inches(0.4), TOP, SLIDE_W - Inches(0.8), Inches(0.5), fill=RGBColor(0xFF, 0xF3, 0xCD))
    text_box(s, Inches(0.6), TOP + Inches(0.07), SLIDE_W - Inches(1.2), Inches(0.38),
             "🔁  Points récurrents sur ce site : §7.1.5 Étalonnage · §8.7 NC — badges dans la check-list",
             size=Pt(10), bold=True, color=C_AMBER)
    TOP += Inches(0.6)

    # 3 colonnes
    col_w = Inches(4.0)
    pad   = Inches(0.15)

    # Col 1 — sélection check-list
    card(s, Inches(0.4), TOP, col_w, Inches(3.8), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(0.15), Inches(3.6), Inches(0.35),
             "1. SÉLECTION DU POINT", size=Pt(9), bold=True, color=C_INK_MUTED)
    items_chk = [
        ("§7.1.5", "Vérification certificats d'étalonnage", C_RED, True),
        ("§7.1.5", "Registre des équipements à jour", C_BRAND, False),
        ("§8.7", "Traçabilité des NC produits", C_AMBER, False),
        ("§8.7", "Zone de quarantaine délimitée", C_AMBER, False),
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.6), Inches(3.6), Inches(3.0))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (cl, label, c, selected) in enumerate(items_chk):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        icon = "●" if selected else "○"
        run.text = f"{icon}  [{cl}] {label}"
        run.font.size  = Pt(10)
        run.font.bold  = selected
        run.font.color.rgb = c if selected else C_INK_MUTED

    # Col 2 — capture
    card(s, Inches(0.4) + col_w + pad, TOP, col_w, Inches(3.8), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6) + col_w + pad, TOP + Inches(0.15), Inches(3.6), Inches(0.35),
             "2. CAPTURE TERRAIN", size=Pt(9), bold=True, color=C_INK_MUTED)
    capture_steps = [
        ("🎙 Micro", "Dicter l'observation à la main libre"),
        ("📷 Photo", "Capturer la preuve visuelle"),
        ("✓/✗ Oui/Non", "Répondre aux 3 questions Claude"),
        ("⚡ Analyser", "Déclencher le diagnostic IA"),
    ]
    txb = s.shapes.add_textbox(Inches(0.6) + col_w + pad, TOP + Inches(0.6), Inches(3.6), Inches(3.0))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (icon, label) in enumerate(capture_steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"{icon}  {label}"
        run.font.size  = Pt(11)
        run.font.color.rgb = C_INK

    # Observation exemple
    rect(s, Inches(4.7), TOP + Inches(2.2), Inches(3.85), Inches(1.1),
         fill=RGBColor(0xF5, 0xF5, 0xF5), line_color=C_DIVIDER, line_w=Pt(0.5))
    text_box(s, Inches(4.85), TOP + Inches(2.3), Inches(3.55), Inches(0.9),
             "« Les clés dynamométriques du poste 12\nne sont pas étalonnées.\nCertificats périmés depuis 8 mois. »",
             size=Pt(10), italic=True, color=C_INK_MUTED)

    # Col 3 — questions oui/non
    card(s, Inches(0.4) + 2 * (col_w + pad), TOP, col_w, Inches(3.8), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6) + 2 * (col_w + pad), TOP + Inches(0.15), Inches(3.6), Inches(0.35),
             "3. RÉPONSES SUGGÉRÉES IA", size=Pt(9), bold=True, color=C_INK_MUTED)
    questions = [
        ("✓ OUI", "Certificats datés < 12 mois ?", C_EMERALD),
        ("✗ NON", "Équipement couvert COFRAC ?", C_RED),
        ("—", "Registre équipements à jour ?", C_INK_MUTED),
    ]
    txb = s.shapes.add_textbox(Inches(0.6) + 2 * (col_w + pad), TOP + Inches(0.6), Inches(3.6), Inches(3.0))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (reponse, question, c) in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"{reponse}"
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = c
        run2 = p.add_run()
        run2.text = f"  {question}"
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK

    text_box(s, Inches(0.6) + 2 * (col_w + pad), TOP + Inches(2.6), Inches(3.6), Inches(0.8),
             "Générées par Claude Opus\n(POST /questions_oui_non)",
             size=Pt(9), italic=True, color=C_INK_MUTED)


def slide_inspection_2(prs):
    """Slide 9 — Acte II : Résultat d'analyse + capitalisation"""
    s = blank(prs)
    header_bar(s, "ACTE II — Résultat de l'analyse IA",
               "Non-conformité majeure · Récidive · Capitalisation de l'expertise")
    page_bg(s)

    TOP = Inches(1.35)

    # Card résultat
    card(s, Inches(0.4), TOP, Inches(7.5), Inches(3.6), fill=C_SURFACE,
         border=C_RED)
    rect(s, Inches(0.4), TOP, Inches(7.5), Inches(0.48), fill=C_RED)
    text_box(s, Inches(0.4), TOP + Inches(0.06), Inches(7.5), Inches(0.38),
             "NC MAJEURE — §7.1.5  Étalonnage et surveillance des équipements de mesure",
             size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    result_items = [
        ("Constat", "Clés dynamométriques poste 12 non étalonnées. Certificats périmés depuis 8 mois."),
        ("Action corrective", "Étalonnage COFRAC immédiat des équipements concernés."),
        ("Action préventive", "Mise en place d'un planning automatisé de suivi des certificats."),
        ("Délai", "Sous 72 h · Responsable : Mei Lin Zhang"),
        ("Source norme", "ISO 9001:2015 §7.1.5.1 — extrait verbatim au survol"),
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.6), Inches(7.1), Inches(2.8))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(result_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{k} : "
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(11)
        run2.font.color.rgb = C_INK

    # Badge récidive
    rect(s, Inches(0.55), TOP + Inches(3.15), Inches(3.2), Inches(0.32), fill=C_AMBER)
    text_box(s, Inches(0.55), TOP + Inches(3.15), Inches(3.2), Inches(0.32),
             "🔁  RÉCIDIVE — NC déjà signalée en nov. 2024",
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Badge Claude
    rect(s, Inches(4.0), TOP + Inches(3.15), Inches(2.2), Inches(0.32), fill=C_BRAND)
    text_box(s, Inches(4.0), TOP + Inches(3.15), Inches(2.2), Inches(0.32),
             "Claude Opus · GEP Capgemini",
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Col droite — capitalisation + RAG
    card(s, Inches(8.2), TOP, Inches(4.9), Inches(3.6), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(8.4), TOP + Inches(0.15), Inches(4.5), Inches(0.35),
             "CAPITALISATION + ARTICLES RAG", size=Pt(9), bold=True, color=C_INK_MUTED)

    capi_items = [
        ("👍 Confirmer", "Marc valide le diagnostic — il entre dans la base"),
        ("✏️  Corriger", "Marc ajuste — sa correction est enregistrée"),
        ("23 ans", "d'expertise BV capturés et redistribuables"),
        ("", ""),
        ("Articles ISO", "Mis à jour au clic sur l'item de check-list"),
        ("Tooltip RAG", "Extrait verbatim de la norme au survol"),
        ("Explicabilité", "Réponse directe à l'exigence CDC"),
    ]
    txb = s.shapes.add_textbox(Inches(8.4), TOP + Inches(0.6), Inches(4.5), Inches(2.9))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(capi_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{k}  " if k else ""
        run.font.size  = Pt(10)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND if k else C_INK
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK_MUTED

    # Douleurs traitées
    rect(s, Inches(0.4), Inches(5.25), SLIDE_W - Inches(0.8), Inches(0.65), fill=C_BRAND)
    text_box(s, Inches(0.6), Inches(5.32), SLIDE_W - Inches(1.2), Inches(0.5),
             "Douleurs traitées :  ③ Historique exploité (🔁 récidive)  ·  ④ Constats objectifs  ·  ① Expertise capitalisée",
             size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 2nd constat
    card(s, Inches(0.4), Inches(6.1), SLIDE_W - Inches(0.8), Inches(0.65),
         fill=C_SURFACE, border=C_AMBER)
    rect(s, Inches(0.4), Inches(6.1), Inches(0.6), Inches(0.65), fill=C_AMBER)
    text_box(s, Inches(1.1), Inches(6.2), SLIDE_W - Inches(1.6), Inches(0.45),
             "2ème constat : NC MINEURE §8.7 — Zone de quarantaine non délimitée, pièces NC mélangées aux pièces valides.",
             size=Pt(10), color=C_INK)


def slide_rapport(prs):
    """Slide 10 — Acte III : Rapport"""
    s = blank(prs)
    header_bar(s, "ACTE III — Rapport instantané",
               "Sur site, en temps réel — plus de soirée de rédaction — T 6:00")
    page_bg(s)

    TOP = Inches(1.35)
    acte_badge(s, Inches(0.4), TOP, "ACTE III", "APRÈS", C_DARK_TEAL)
    TOP += Inches(0.5)

    # Col gauche
    card(s, Inches(0.4), TOP, Inches(4.1), Inches(4.5), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(0.15), Inches(3.7), Inches(0.35),
             "SYNTHÈSE DU RAPPORT", size=Pt(9), bold=True, color=C_INK_MUTED)

    synth_items = [
        ("NC MAJEURES", "2", C_RED),
        ("NC MINEURES", "0", C_AMBER),
        ("OBSERVATIONS", "0", C_BRAND),
        ("CONFORMES", "0", C_EMERALD),
    ]
    for i, (label, val, c) in enumerate(synth_items):
        cy = TOP + Inches(0.65) + i * Inches(0.48)
        rect(s, Inches(0.6), cy, Inches(0.45), Inches(0.38), fill=c)
        text_box(s, Inches(0.6), cy + Inches(0.05), Inches(0.45), Inches(0.3),
                 val, size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        text_box(s, Inches(1.15), cy + Inches(0.07), Inches(3.1), Inches(0.3),
                 label, size=Pt(11), color=C_INK)

    # Grille conformité
    text_box(s, Inches(0.6), TOP + Inches(2.7), Inches(3.7), Inches(0.35),
             "GRILLE DE CONFORMITÉ", size=Pt(9), bold=True, color=C_INK_MUTED)
    conformite = [
        ("§7.1.5 Étalonnage", 18, C_RED),
        ("§8.7 Non-conformités", 42, C_AMBER),
        ("§7.5 Documentation", 80, C_EMERALD),
        ("§7.2 Compétences", 85, C_EMERALD),
    ]
    for i, (label, pct, c) in enumerate(conformite):
        cy = TOP + Inches(3.1) + i * Inches(0.3)
        text_box(s, Inches(0.6), cy, Inches(2.0), Inches(0.25),
                 label, size=Pt(9), color=C_INK_MUTED)
        bar_total = Inches(1.5)
        rect(s, Inches(2.7), cy + Inches(0.04), bar_total, Inches(0.17), fill=C_DIVIDER)
        rect(s, Inches(2.7), cy + Inches(0.04), int(bar_total * pct / 100), Inches(0.17), fill=c)
        text_box(s, Inches(4.3), cy, Inches(0.6), Inches(0.25),
                 f"{pct}%", size=Pt(9), bold=True, color=c)

    # Col centre
    card(s, Inches(4.75), TOP, Inches(4.2), Inches(4.5), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(4.95), TOP + Inches(0.15), Inches(3.8), Inches(0.35),
             "ACTIONS DISPONIBLES", size=Pt(9), bold=True, color=C_INK_MUTED)
    actions = [
        ("✏️  Modifier", "Marc relit et ajuste une formulation — l'IA propose, l'auditeur décide"),
        ("🔒 Anonymiser", "RGPD — les noms disparaissent de l'aperçu en un clic"),
        ("⬇ Télécharger", "Export Word (.docx) généré par le backend"),
        ("✍️ Signer", "Validation contradictoire — signature manuscrite Mei Lin Zhang"),
        ("⬆ Portail RATP", "Envoi du rapport signé → portail fournisseur RATP"),
    ]
    txb = s.shapes.add_textbox(Inches(4.95), TOP + Inches(0.6), Inches(3.8), Inches(3.7))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (action, desc) in enumerate(actions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"{action}  "
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = desc
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK_MUTED

    # Col droite — le geste clé
    card(s, Inches(9.2), TOP, Inches(3.9), Inches(4.5), fill=RGBColor(0xE8, 0xF5, 0xE9), border=C_EMERALD)
    rect(s, Inches(9.2), TOP, Inches(3.9), Inches(0.45), fill=C_DARK_TEAL)
    text_box(s, Inches(9.2), TOP + Inches(0.05), Inches(3.9), Inches(0.37),
             "⬆  LE GESTE QUI FERME LA BOUCLE",
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    fermer_items = [
        "Clic sur « ⬆ Portail RATP »",
        "→ Génère le .docx (backend)",
        "→ Encode en base64",
        "→ Crée un document fromBV",
        "→ Stocke dans localStorage",
        "",
        "Le bouton passe à :",
        "« ✓ Transmis au portail »",
        "",
        "Marc quitte le site.",
        "Le rapport est déjà chez RATP.",
    ]
    txb = s.shapes.add_textbox(Inches(9.4), TOP + Inches(0.6), Inches(3.5), Inches(3.7))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(fermer_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(10)
        run.font.bold  = item.startswith("«")
        run.font.color.rgb = C_BRAND if item.startswith("→") else (C_DARK_TEAL if item.startswith("«") else C_INK)

    # Douleur 5
    rect(s, Inches(0.4), Inches(6.05), SLIDE_W - Inches(0.8), Inches(0.65), fill=C_DARK_TEAL)
    text_box(s, Inches(0.6), Inches(6.12), SLIDE_W - Inches(1.2), Inches(0.5),
             "Douleur traitée :  ⑤ Rapport généré sur place — la soirée de rédaction a disparu",
             size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_epilogue(prs):
    """Slide 11 — Épilogue : Portail RATP"""
    s = blank(prs)
    header_bar(s, "ÉPILOGUE — La boucle se referme",
               "Portail fournisseur RATP · Mei Lin Zhang — T 7:00")
    page_bg(s)

    TOP = Inches(1.35)

    # Titre avec accent
    rect(s, Inches(0.4), TOP, SLIDE_W - Inches(0.8), Inches(0.52), fill=C_EMERALD)
    text_box(s, Inches(0.6), TOP + Inches(0.07), SLIDE_W - Inches(1.2), Inches(0.4),
             "Mei Lin Zhang se connecte — et trouve le rapport de Marc",
             size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    TOP += Inches(0.65)

    # Col gauche — avant
    card(s, Inches(0.4), TOP, Inches(5.9), Inches(3.8), fill=C_SURFACE, border=C_DIVIDER)
    text_box(s, Inches(0.6), TOP + Inches(0.15), Inches(5.5), Inches(0.35),
             "AVANT L'AUDIT — Documents déposés par Mei Lin Zhang",
             size=Pt(9), bold=True, color=C_INK_MUTED)
    docs_avant = [
        ("✓  CR_Audit_nov_2024.docx", "Compte-rendu · pré-analysé par Claude"),
        ("✓  Procedures_etalonnage_v3.pdf", "Procédure interne · équipements périmés"),
        ("✓  Plan_qualite_2025.pdf", "Plan qualité · budget non alloué"),
    ]
    txb = s.shapes.add_textbox(Inches(0.6), TOP + Inches(0.6), Inches(5.5), Inches(2.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (nom, note) in enumerate(docs_avant):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = nom
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = C_EMERALD
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = f"  {note}"
        run2.font.size  = Pt(10)
        run2.font.color.rgb = C_INK_MUTED

    text_box(s, Inches(0.6), TOP + Inches(3.2), Inches(5.5), Inches(0.45),
             "Catégories : Audit & Inspection · Procédures · Sécurité & Risques",
             size=Pt(9), italic=True, color=C_INK_MUTED)

    # Flèche centrale
    text_box(s, Inches(6.5), TOP + Inches(1.5), Inches(0.8), Inches(0.8),
             "→", size=Pt(36), bold=True, color=C_BRAND, align=PP_ALIGN.CENTER)

    # Col droite — rapport BV
    card(s, Inches(7.2), TOP, Inches(5.7), Inches(3.8), fill=RGBColor(0xE3, 0xF2, 0xFD), border=C_BRAND)
    rect(s, Inches(7.2), TOP, Inches(5.7), Inches(0.48), fill=C_BRAND)
    text_box(s, Inches(7.2), TOP + Inches(0.06), Inches(5.7), Inches(0.38),
             "NOUVEAU  ·  Rapport d'audit  ·  Badge Bureau Veritas",
             size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rapport_items = [
        ("Type", "Rapport d'audit BV · Badge bleu Bureau Veritas"),
        ("Analyse", "2 NC majeures · clauses §7.1.5 et §8.7"),
        ("Actions", "Étalonnage COFRAC < 72 h · Zone quarantaine"),
        ("Fichier", "⬇ Télécharger le .docx signé — disponible"),
        ("Persistance", "localStorage — survit à F5 + changement de compte"),
    ]
    txb = s.shapes.add_textbox(Inches(7.4), TOP + Inches(0.65), Inches(5.3), Inches(3.0))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (k, v) in enumerate(rapport_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"{k} : "
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = C_BRAND
        run2 = p.add_run()
        run2.text = v
        run2.font.size  = Pt(11)
        run2.font.color.rgb = C_INK

    # Bas — the payoff
    rect(s, Inches(0.4), Inches(5.3), SLIDE_W - Inches(0.8), Inches(1.0), fill=C_DARK_TEAL)
    text_box(s, Inches(0.6), Inches(5.42), SLIDE_W - Inches(1.2), Inches(0.4),
             "RATP dépose  →  BV audite  →  le rapport revient à RATP",
             size=Pt(16), bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
    text_box(s, Inches(0.6), Inches(5.85), SLIDE_W - Inches(1.2), Inches(0.4),
             "Une seule plateforme · Un seul moteur IA · La boucle est complète",
             size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)


def slide_cloture(prs):
    """Slide 12 — Clôture : Les 5 valeurs"""
    s = blank(prs)

    # Fond dark-teal
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK_TEAL)

    # Bandeau brand bas
    rect(s, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), fill=C_BRAND)

    # Titre
    text_box(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Les cinq promesses — toutes tenues",
             size=Pt(28), bold=True, color=C_WHITE)
    text_box(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
             "UC 28 — Inspection Augmentée · Code Resonance · Hackathon Capgemini × Anthropic 2026",
             size=Pt(11), color=C_CYAN, italic=True)

    # Ligne de séparation
    rect(s, Inches(0.5), Inches(1.4), Inches(5), Inches(0.04), fill=C_CYAN)

    # 5 valeurs
    valeurs = [
        ("①", "Standardisation", C_BRAND,
         "Check-list générée — chaque auditeur part avec la même base."),
        ("②", "Gain de temps", C_CYAN,
         "Brief en 30 s, rapport sur place. La soirée de rédaction : supprimée."),
        ("③", "Respect des normes", C_EMERALD,
         "Chaque constat sourcé — l'article exact ISO à l'écran, verbatim."),
        ("④", "Capitalisation expertise", C_AMBER,
         "Confirmer/Corriger — 23 ans d'expertise BV capturés et redistribués."),
        ("⑤", "Qualité & traçabilité", RGBColor(0xAB, 0x47, 0xBC),
         "Tout est daté, sourcé, signé. Explicabilité native — exigence CDC."),
    ]
    for i, (num, titre, c, desc) in enumerate(valeurs):
        cy = Inches(1.6) + i * Inches(1.0)
        rect(s, Inches(0.5), cy, Inches(0.5), Inches(0.6), fill=c)
        text_box(s, Inches(0.5), cy + Inches(0.1), Inches(0.5), Inches(0.4),
                 num, size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        text_box(s, Inches(1.15), cy + Inches(0.04), Inches(3.5), Inches(0.3),
                 titre, size=Pt(13), bold=True, color=c)
        text_box(s, Inches(1.15), cy + Inches(0.32), Inches(7.5), Inches(0.35),
                 desc, size=Pt(11), color=C_WHITE)

    # Séparateur vertical
    rect(s, Inches(9.3), Inches(1.4), Inches(0.04), Inches(5.5), fill=C_CYAN)

    # Impact
    rect(s, Inches(9.5), Inches(1.5), Inches(3.5), Inches(2.4), fill=RGBColor(0x00, 0x4A, 0x68))
    text_box(s, Inches(9.65), Inches(1.65), Inches(3.2), Inches(0.4),
             "90% PAPERASSE", size=Pt(22), bold=True,
             color=RGBColor(0xFF, 0x6B, 0x6B), align=PP_ALIGN.CENTER)
    text_box(s, Inches(9.65), Inches(2.1), Inches(3.2), Inches(0.35),
             "10% TERRAIN", size=Pt(18), bold=True, color=C_INK_MUTED,
             align=PP_ALIGN.CENTER)
    text_box(s, Inches(9.65), Inches(2.5), Inches(3.2), Inches(0.35),
             "↕  INVERSÉ", size=Pt(20), bold=True, color=C_CYAN,
             align=PP_ALIGN.CENTER)
    text_box(s, Inches(9.65), Inches(2.9), Inches(3.2), Inches(0.35),
             "10% PAPERASSE", size=Pt(14), color=C_EMERALD,
             align=PP_ALIGN.CENTER, bold=True)
    text_box(s, Inches(9.65), Inches(3.28), Inches(3.2), Inches(0.35),
             "90% TERRAIN", size=Pt(20), bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    # KPI ≈ 24 000 h/an
    rect(s, Inches(9.5), Inches(4.1), Inches(3.5), Inches(1.1), fill=C_BRAND)
    text_box(s, Inches(9.5), Inches(4.18), Inches(3.5), Inches(0.5),
             "≈ 24 000 h / an", size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text_box(s, Inches(9.5), Inches(4.68), Inches(3.5), Inches(0.4),
             "potentiel récupéré", size=Pt(11), color=C_CYAN, align=PP_ALIGN.CENTER)

    # Clôture
    text_box(s, Inches(0.5), Inches(6.65), Inches(8.5), Inches(0.45),
             "Merci.  · Philippe Gragnic · Habib Koffi · Véronique Poilane Zhang",
             size=Pt(12), color=C_WHITE)
    text_box(s, Inches(9.5), Inches(6.65), Inches(3.5), Inches(0.45),
             "Équipe Code Resonance", size=Pt(11), color=C_CYAN,
             align=PP_ALIGN.RIGHT)


# ── Slide screenshot ─────────────────────────────────────────────────────────

def slide_screenshot(prs, img_filename, title, annotations, accent_color=None):
    """Slide visuel : screenshot pleine page + bandeau titre haut + annotations bas."""
    s    = blank(prs)
    ac   = accent_color or C_BRAND
    img_path = os.path.abspath(os.path.join(SCREENSHOTS, img_filename))
    if not os.path.exists(img_path):
        return  # skip gracefully if file missing

    # Fond noir pour contraste
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_DARK_TEAL)

    # Bandeau titre haut (54px)
    rect(s, 0, 0, SLIDE_W, Inches(0.72), fill=ac)
    text_box(s, Inches(0.3), Inches(0.1), Inches(10), Inches(0.52),
             title, size=Pt(15), bold=True, color=C_WHITE)
    # Numéro slide discret
    text_box(s, SLIDE_W - Inches(1.2), Inches(0.18), Inches(1.0), Inches(0.38),
             "ÉCRAN →", size=Pt(9), color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.RIGHT, italic=True)

    # Screenshot (centré, hauteur = slide - titre - bande bas)
    BAND_H = Inches(0.82)
    img_top    = Inches(0.78)
    img_height = SLIDE_H - img_top - BAND_H
    # Aspect 1280/720 = 16/9 → largeur max = 13.33 in mais marges 0.2
    img_width  = img_height * (1280 / 720)
    if img_width > SLIDE_W - Inches(0.4):
        img_width  = SLIDE_W - Inches(0.4)
        img_height = img_width * (720 / 1280)
    img_left = (SLIDE_W - img_width) / 2

    s.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)

    # Bandeau annotations bas
    band_top = SLIDE_H - BAND_H
    rect(s, 0, band_top, SLIDE_W, BAND_H, fill=C_INK)
    if annotations:
        col_w = SLIDE_W / len(annotations)
        for i, (icon, txt) in enumerate(annotations):
            cx = i * col_w
            # Séparateur vertical entre colonnes
            if i > 0:
                rect(s, cx, band_top + Inches(0.1), Inches(0.015),
                     BAND_H - Inches(0.2), fill=RGBColor(0x40, 0x55, 0x60))
            text_box(s, cx + Inches(0.2), band_top + Inches(0.1),
                     col_w - Inches(0.3), Inches(0.35),
                     icon, size=Pt(10), bold=True, color=ac)
            text_box(s, cx + Inches(0.2), band_top + Inches(0.44),
                     col_w - Inches(0.3), Inches(0.32),
                     txt, size=Pt(9), color=RGBColor(0xCC, 0xD8, 0xE0))


# ── Slide de transition (mini slide entre actes) ───────────────────────────

def slide_transition(prs, acte, label, color, subtitle=""):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=color)

    text_box(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
             acte, size=Pt(60), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    text_box(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.7),
             label, size=Pt(28), color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)
    if subtitle:
        text_box(s, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.45),
                 subtitle, size=Pt(14), color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_cover(prs)            # 1
    slide_douleurs(prs)         # 2
    slide_architecture(prs)     # 3
    slide_connexion(prs)        # 4

    slide_dashboard(prs)        # 5
    slide_screenshot(prs, "01_dashboard.png",
        "Écran 1 — Dashboard · La journée de Marc d'un seul regard",
        [("10 missions planifiées", "Trait rouge temps réel · vue Liste/Planning"),
         ("Carte Leaflet + OSRM", "Itinéraire routier réel, clic → zoom sur le site"),
         ("Filtre PROCHAIN", "Mission Sucy-en-Brie sélectionnée → Démarrer"),
         ("Hors-ligne", "Zéro clé API — fonctionne même sans réseau")],
        C_BRAND)                # 6

    slide_transition(prs,        # 7
        "ACTE I", "Avant l'audit — Le Brief",
        C_BRAND, "T 2:20 · Marc n'a pas quitté son bureau — son brief est déjà prêt")

    slide_brief_1(prs)          # 8
    slide_screenshot(prs, "02_brief_context.png",
        "ACTE I · Brief — Contexte site, historique NC, check-list auto-générée",
        [("Historique branché sur audit.db", "NC §7.1.5 ouverte depuis nov. 2024 — jamais clôturée"),
         ("Thèmes récurrents 🔁", "§7.1.5 + §8.7 → badges dans la check-list d'inspection"),
         ("4 sections générées", "S1 §7.5 · S2 §7.1.5 · S3 §8.7 · S4 §7.2"),
         ("Douleur ③ traitée", "Historique exploité — Marc n'a rien eu à se rappeler")],
        C_BRAND)                # 9

    slide_brief_2(prs)          # 10
    slide_screenshot(prs, "03_brief_scope.png",
        "ACTE I · Brief — Scope, documents RATP, check-list éditable",
        [("Scope de l'audit", "3 domaines actifs · §7.1.5 · §8.7 · §7.5"),
         ("⚠ Alertes RATP", "3 docs pré-analysés par Claude · risques remontés"),
         ("Check-list éditable", "× supprimer un item · + ajouter un point sur mesure"),
         ("Douleurs ① ②", "Expertise capitalisée (badge Auditeur) · Prépa standardisée")],
        C_BRAND)                # 11

    slide_transition(prs,        # 12
        "ACTE II", "Pendant l'audit — Le Terrain",
        C_EMERALD, "T 3:30 · Marc est sur site, gants aux mains, devant les équipements")

    slide_inspection_1(prs)     # 13
    slide_screenshot(prs, "04_inspection.png",
        "ACTE II · Inspection — Capture terrain, analyse IA, questions oui/non",
        [("Bannière 🔁 récurrences", "§7.1.5 + §8.7 rappelés avant même de commencer"),
         ("NC MAJEURE §7.1.5", "Clés dynamométriques non étalonnées · certificats périmés 8 mois"),
         ("3 questions Claude", "Oui/Non contextualisées · réponses réinjectées dans l'analyse"),
         ("Douleurs ③ ④ traitées", "Récidive détectée · constat objectif généré par l'IA")],
        C_EMERALD)              # 14

    slide_inspection_2(prs)     # 15

    slide_transition(prs,        # 16
        "ACTE III", "Après l'audit — Le Rapport",
        C_DARK_TEAL, "T 6:00 · Le rapport est déjà là. Sur place. Pas ce soir — maintenant.")

    slide_rapport(prs)          # 17
    slide_screenshot(prs, "05_report.png",
        "ACTE III · Rapport — Synthèse, grille conformité, signature, envoi portail",
        [("Score 35% — alerte rouge", "§7.1.5 à 18% · §8.7 à 42% · barres colorées par seuil"),
         ("Actions disponibles", "Modifier · Anonymiser RGPD · Télécharger Word · Signer"),
         ("⬆ Portail RATP", "Un clic → rapport transmis directement chez Mei Lin Zhang"),
         ("Douleur ⑤ traitée", "La soirée de rédaction de Marc a disparu")],
        C_DARK_TEAL)            # 18

    slide_transition(prs,        # 19
        "ÉPILOGUE", "La boucle se referme",
        RGBColor(0x00, 0x5C, 0x40), "T 7:00 · Mei Lin Zhang se connecte au portail RATP")

    slide_epilogue(prs)         # 20
    slide_screenshot(prs, "06_portal.png",
        "ÉPILOGUE · Portail RATP — Mei Lin Zhang reçoit le rapport de Marc",
        [("Badge Bureau Veritas", "Le rapport de Marc apparaît en tête de liste"),
         ("Analyse instantanée", "2 NC majeures · actions · clauses à risque déjà visibles"),
         ("⬇ Télécharger", "Le .docx signé est récupérable en un clic"),
         ("Boucle complète", "RATP dépose → BV audite → rapport revient à RATP")],
        RGBColor(0x00, 0x5C, 0x40))  # 21

    slide_cloture(prs)          # 22

    out = os.path.abspath(OUT_PATH)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"✓ Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
