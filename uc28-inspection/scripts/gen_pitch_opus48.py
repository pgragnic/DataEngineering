"""
Génère UC28_Pitch_Deck_Opus48.pptx — version 3 actes (script Opus 4.8)
Palette : Capgemini #0070AD + émeraude #059669 + ardoise sombre
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ────────────────────────────────────────────────────────────────
BRAND    = RGBColor(0x00, 0x70, 0xAD)   # Capgemini bleu
EMERALD  = RGBColor(0x05, 0x96, 0x69)   # vert émeraude
SLATE    = RGBColor(0x1E, 0x29, 0x3B)   # ardoise sombre (titres)
MIST     = RGBColor(0xEE, 0xF2, 0xF7)   # fond app
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x64, 0x74, 0x8B)   # texte secondaire
RED_NC   = RGBColor(0xDC, 0x26, 0x26)   # NC majeure
AMBER    = RGBColor(0xD9, 0x77, 0x06)   # NC mineure
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)   # fond slides contenu

# ── Dimensions 16:9 ────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # layout vide


# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box


def slide_header(slide, acte_label, title, subtitle=None,
                 accent=BRAND, bg=LIGHT_BG):
    """Bande de fond + bandeau haut coloré + titres."""
    add_rect(slide, 0, 0, W, H, bg)
    # bandeau haut
    add_rect(slide, 0, 0, W, Inches(1.6), accent)
    # étiquette acte (petite caps)
    if acte_label:
        add_text(slide, acte_label,
                 Inches(0.5), Inches(0.15), Inches(12), Inches(0.45),
                 size=11, bold=True, color=WHITE)
    # titre principal
    add_text(slide, title,
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.85),
             size=34, bold=True, color=WHITE)
    # sous-titre
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
                 size=16, bold=False, color=RGBColor(0xBC, 0xD9, 0xED))


def bullet(slide, items, x, y, w,
           size=18, color=SLATE, spacing=Inches(0.38),
           dot_color=BRAND, dot="▸"):
    """Liste à puces simple."""
    for i, item in enumerate(items):
        add_text(slide, f"{dot}  {item}",
                 x, y + i * spacing, w, spacing,
                 size=size, color=color)


def pill(slide, text, x, y, w=Inches(2.8), h=Inches(0.42),
         bg=BRAND, fg=WHITE, size=14):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, text, x, y + Inches(0.04), w, h,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def card(slide, title, body_lines, x, y, w=Inches(3.8), h=Inches(2.2),
         title_color=BRAND, bg=WHITE, border=BRAND):
    add_rect(slide, x, y, w, h, bg)
    # bordure gauche
    add_rect(slide, x, y, Inches(0.06), h, border)
    add_text(slide, title,
             x + Inches(0.18), y + Inches(0.12), w - Inches(0.25), Inches(0.38),
             size=13, bold=True, color=title_color)
    for i, line in enumerate(body_lines):
        add_text(slide, line,
                 x + Inches(0.18), y + Inches(0.55) + i * Inches(0.34),
                 w - Inches(0.25), Inches(0.34),
                 size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — L'ACCROCHE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, W, H, SLATE)

# grande citation
add_text(s1, "Il est 21 heures.",
         Inches(1), Inches(1.2), Inches(11.3), Inches(0.7),
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1,
         "Marc a terminé son audit à 16 heures.\n"
         "Et il est encore devant son écran.",
         Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.1),
         size=24, color=RGBColor(0xBC, 0xD9, 0xED), align=PP_ALIGN.CENTER)

# stat 90/10
add_rect(s1, Inches(4.2), Inches(3.4), Inches(4.9), Inches(1.05), BRAND)
add_text(s1, "90 % paperasse · 10 % terrain",
         Inches(4.2), Inches(3.48), Inches(4.9), Inches(0.9),
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1, "Nous, on a inversé ce ratio.",
         Inches(1), Inches(4.7), Inches(11.3), Inches(0.65),
         size=30, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

add_text(s1, "UC 28 — Inspection Augmentée · Code Resonance · Hackathon Capgemini × Anthropic 2026",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CADRE : ÉQUIPE, MODÈLE, 5 DOULEURS
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_header(s2, "LE CADRE", "Équipe · Modèle · Enjeu")

# Équipe
add_text(s2, "Code Resonance",
         Inches(0.5), Inches(1.9), Inches(4), Inches(0.45),
         size=16, bold=True, color=BRAND)
bullet(s2,
       ["Habib KOFFI — Agent IA",
        "Véronique POILANE ZHANG — UI/UX",
        "Philippe GRAGNIC — Data & démo"],
       Inches(0.5), Inches(2.4), Inches(4),
       size=14, color=SLATE, spacing=Inches(0.35), dot="·")

# Modèle économique
add_rect(s2, Inches(5), Inches(1.85), Inches(7.8), Inches(1.45), MIST)
add_text(s2, "Modèle B2B",
         Inches(5.2), Inches(1.92), Inches(7.4), Inches(0.4),
         size=13, bold=True, color=BRAND)
add_text(s2,
         "Bureau Veritas  achète  →  audite  RATP  et  Apave\n"
         "Une plateforme, deux faces — un seul moteur d'IA",
         Inches(5.2), Inches(2.32), Inches(7.4), Inches(0.85),
         size=14, color=SLATE)

# 5 douleurs
add_text(s2, "Les 5 douleurs adressées (CDC officiel)",
         Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=SLATE)

douleurs = [
    "Expertise rare et difficile à transmettre entre inspecteurs",
    "Préparation longue et hétérogène selon les auditeurs",
    "Historique des NC sous-exploité sur le terrain",
    "Constats parfois subjectifs ou incomplets",
    "Rédaction des rapports chronophage (soirées entières)",
]
for i, d in enumerate(douleurs):
    add_text(s2, f"{i+1}.  {d}",
             Inches(0.7), Inches(3.95) + i * Inches(0.5), Inches(12), Inches(0.46),
             size=14, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ACTE I · AVANT · LE BRIEF
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
slide_header(s3, "ACTE I — AVANT L'AUDIT", "Le Brief — préparation en 30 secondes",
             subtitle="Marc n'a pas encore quitté son bureau. Son brief est déjà prêt.")

# Ligne du temps
add_text(s3, "CONNEXION", Inches(0.4), Inches(1.85), Inches(1.5), Inches(0.35),
         size=10, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
add_text(s3, "→", Inches(1.9), Inches(1.85), Inches(0.4), Inches(0.35),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s3, "CLIENT\nRATP", Inches(2.3), Inches(1.75), Inches(1.2), Inches(0.55),
         size=10, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
add_text(s3, "→", Inches(3.5), Inches(1.85), Inches(0.4), Inches(0.35),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s3, "DASHBOARD", Inches(3.9), Inches(1.85), Inches(1.5), Inches(0.35),
         size=10, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
add_text(s3, "→", Inches(5.4), Inches(1.85), Inches(0.4), Inches(0.35),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s3, Inches(5.8), Inches(1.75), Inches(1.5), Inches(0.5), BRAND)
add_text(s3, "BRIEF ★", Inches(5.8), Inches(1.8), Inches(1.5), Inches(0.4),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 cartes
card(s3, "Contexte site",
     ["Atelier Sucy-en-Brie — RER A",
      "218 personnes",
      "Rames MI09 · 2 h 30"],
     Inches(0.4), Inches(2.5))

card(s3, "⚠  Alerte historique",
     ["NC mineure §7.1.5",
      "Ouverte : nov. 2024",
      "Non clôturée → priorité S2"],
     Inches(4.5), Inches(2.5),
     title_color=RED_NC, border=RED_NC)

card(s3, "Check-list auto-générée",
     ["5 sources : référentiel, scope,",
      "contexte site, historique NC,",
      "docs fournisseur pré-analysés"],
     Inches(8.6), Inches(2.5),
     title_color=EMERALD, border=EMERALD)

# Douleur résolue
pill(s3, "✓  Douleur 1 & 2 : historique exploité · prépa standardisée",
     Inches(0.5), Inches(5.05), w=Inches(12.3), h=Inches(0.48),
     bg=EMERALD)

# Note actions
add_text(s3,
         "Marc retire les items non pertinents (×) · ajoute ses propres points (badge Auditeur) ·"
         " les alertes RATP remontent depuis le portail fournisseur",
         Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.55),
         size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ACTE II · PENDANT · LE TERRAIN
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
slide_header(s4, "ACTE II — PENDANT L'AUDIT", "Le Terrain — diagnostic en temps réel",
             subtitle="Marc est sur site. Un seul clic déclenche trois actions simultanées.",
             accent=EMERALD)

# 3 colonnes
col_w = Inches(3.9)
col_y = Inches(1.85)
col_h = Inches(3.6)

# Col 1 — Check-list
add_rect(s4, Inches(0.3), col_y, col_w, col_h, MIST)
add_text(s4, "1. CHECK-LIST", Inches(0.45), col_y + Inches(0.1), col_w, Inches(0.35),
         size=12, bold=True, color=BRAND)
bullet(s4,
       ["Sélection §7.1.5 Étalonnage",
        "→ suggestions générées",
        "→ questions oui/non Claude",
        "→ articles ISO mis à jour",
        "",
        "Micro 🎤 · Photo 📷",
        "Manuscrit ✍ (stylet)"],
       Inches(0.45), col_y + Inches(0.5), col_w - Inches(0.3),
       size=12, color=SLATE, spacing=Inches(0.37), dot="·")

# Col 2 — Diagnostic
add_rect(s4, Inches(4.6), col_y, col_w, col_h, MIST)
add_text(s4, "2. DIAGNOSTIC CLAUDE", Inches(4.75), col_y + Inches(0.1), col_w, Inches(0.35),
         size=12, bold=True, color=BRAND)
add_rect(s4, Inches(4.6), col_y + Inches(0.5), col_w, Inches(0.55), RED_NC)
add_text(s4, "NC MAJEURE · §7.1.5",
         Inches(4.75), col_y + Inches(0.55), col_w, Inches(0.42),
         size=14, bold=True, color=WHITE)
bullet(s4,
       ["Contexte : Sucy, MI09, NC nov. 2024",
        "Action corrective : étalonnage COFRAC",
        "Action préventive : planning auto",
        "🔁 Récidive détectée",
        "Clause sourcée · verbatim à l'écran"],
       Inches(4.75), col_y + Inches(1.15), col_w - Inches(0.3),
       size=12, color=SLATE, spacing=Inches(0.38), dot="·")

# Col 3 — Capitalisation
add_rect(s4, Inches(8.9), col_y, col_w, col_h, MIST)
add_text(s4, "3. CAPITALISATION BV", Inches(9.05), col_y + Inches(0.1), col_w, Inches(0.35),
         size=12, bold=True, color=EMERALD)
bullet(s4,
       ["Bouton Confirmer → expertise validée",
        "Bouton Corriger → donnée d'apprentissage",
        "",
        "23 ans de terrain de Marc",
        "→ capturés, réutilisables",
        "→ transmis à toute l'équipe BV"],
       Inches(9.05), col_y + Inches(0.5), col_w - Inches(0.3),
       size=12, color=SLATE, spacing=Inches(0.37), dot="·")

# Douleurs résolues
pill(s4, "✓  Douleurs 3, 4 & 5 : objectivité · historique · expertise capitalisée",
     Inches(0.5), Inches(5.6), w=Inches(12.3), h=Inches(0.48),
     bg=BRAND)

add_text(s4,
         "L'IA est explicable : chaque diagnostic est sourcé par clause ISO. "
         "Articles verbatim accessibles au survol — exigence réglementaire du CDC.",
         Inches(0.5), Inches(6.22), Inches(12.3), Inches(0.5),
         size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ACTE III · APRÈS · LE RAPPORT
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
slide_header(s5, "ACTE III — APRÈS L'AUDIT", "Le Rapport — sur place, immédiatement",
             subtitle="Marc clique Rapport. Il ne rentrera pas avec une soirée de rédaction.",
             accent=SLATE)

# 4 fonctionnalités
feats = [
    ("Rapport instantané",
     ["Synthèse automatique", "Plan d'actions + délais", "Responsables identifiés"],
     Inches(0.4), BRAND),
    ("Grille de conformité",
     ["Score global calculé", "Barres ISO par section", "Rouge < 50 % — §7.1.5 alerte"],
     Inches(3.55), EMERALD),
    ("Contrôle auditeur",
     ["Modification en ligne", "Anonymisation RGPD 1 clic", "Export Word .docx"],
     Inches(6.7), BRAND),
    ("Signature contradictoire",
     ["Canvas manuscrit (tablette)", "Karim signe sur site", "Rapport tracé & validé"],
     Inches(9.85), SLATE),
]

for title, items, x, col in feats:
    card(s5, title, items, x, Inches(2.1), w=Inches(3.0), h=Inches(2.6),
         title_color=col, border=col)

# Citation finale
add_rect(s5, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.65), MIST)
add_text(s5,
         '"Le rapport est généré. Immédiatement. Sur place. Pas ce soir — maintenant."',
         Inches(0.7), Inches(5.08), Inches(12), Inches(0.5),
         size=15, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

pill(s5, "✓  Douleur 5 réglée : rédaction éliminée · Marc quitte le site avec le travail terminé",
     Inches(0.5), Inches(5.85), w=Inches(12.3), h=Inches(0.48),
     bg=EMERALD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LES 5 VALEURS · CLÔTURE
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
slide_header(s6, "CLÔTURE", "Les 5 valeurs livrées — conformes au cahier des charges")

valeurs = [
    (BRAND,   "Standardisation",
              "Chaque auditeur BV part avec la même check-list générée par l'IA."),
    (EMERALD, "Gain de temps",
              "Brief en 30 s. Rapport sur place. Le 90/10 est inversé."),
    (BRAND,   "Respect des normes",
              "ISO 9001:2015 indexé par RAG — articles verbatim à l'écran."),
    (EMERALD, "Capitalisation expertise BV",
              "Confirmer / Corriger : l'expérience terrain devient réutilisable."),
    (BRAND,   "Qualité · Traçabilité · Explicabilité",
              "Tout est daté, sourcé par clause, signé contradictoirement."),
]

for i, (color, titre, desc) in enumerate(valeurs):
    y = Inches(1.85) + i * Inches(0.98)
    add_rect(s6, Inches(0.5), y, Inches(0.08), Inches(0.72), color)
    add_text(s6, titre,
             Inches(0.78), y, Inches(12), Inches(0.38),
             size=16, bold=True, color=color)
    add_text(s6, desc,
             Inches(0.78), y + Inches(0.35), Inches(12), Inches(0.38),
             size=13, color=GRAY)

# Boucle narrative finale
add_rect(s6, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55), SLATE)
add_text(s6,
         "On a commencé à 21 h, Marc devant son écran. "
         "Désormais, à 16 h, il a fini.  ·  Le 90 / 10 est inversé.",
         Inches(0.7), Inches(6.82), Inches(12), Inches(0.42),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BONUS : PORTAIL FOURNISSEUR
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
slide_header(s7, "[BONUS] LA PLATEFORME À DEUX FACES",
             "Portail fournisseur — Mei Lin Zhang côté RATP",
             subtitle="BV achète la solution · RATP & Apave en sont les bénéficiaires")

add_text(s7, "Une seule IA — deux usages", Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.45),
         size=18, bold=True, color=SLATE)

card(s7, "Côté auditeur BV — Marc",
     ["Check-list alertes RATP intégrées",
      "Diagnostic Claude contextualisé",
      "Rapport sur place"],
     Inches(0.5), Inches(2.6), w=Inches(5.8), h=Inches(2.2),
     title_color=BRAND, border=BRAND)

add_text(s7, "⟷", Inches(6.5), Inches(3.3), Inches(0.8), Inches(0.7),
         size=30, color=BRAND, align=PP_ALIGN.CENTER)

card(s7, "Côté fournisseur RATP — Mei Lin Zhang",
     ["Dépôt de documents (6 catégories)",
      "Analyse Claude → résumé + risques",
      "Alertes remontent dans check-list BV"],
     Inches(7.2), Inches(2.6), w=Inches(5.6), h=Inches(2.2),
     title_color=EMERALD, border=EMERALD)

pill(s7, "Le même moteur croise les données des deux côtés — modèle B2B natif",
     Inches(0.5), Inches(5.1), w=Inches(12.3), h=Inches(0.48), bg=BRAND)

add_text(s7, "Référents MU : JC Formosa (Apave) · Jérôme Fouquet (BV) · Marvin Dufresne (RATP)",
         Inches(0.5), Inches(5.88), Inches(12.3), Inches(0.35),
         size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ÉQUIPE
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, W, H, SLATE)

add_text(s8, "Code Resonance",
         Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s8, "UC 28 — Inspection Augmentée · Hackathon Capgemini × Anthropic 2026",
         Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.4),
         size=15, color=RGBColor(0xBC, 0xD9, 0xED), align=PP_ALIGN.CENTER)

membres = [
    ("Habib KOFFI",           "Agent IA & Backend",   "habib.koffi@capgemini.com"),
    ("Véronique POILANE ZHANG", "UI/UX & Design",     "veronique.poilane-zhang-ext@capgemini.com"),
    ("Philippe GRAGNIC",      "Data & Démo",          "philippe.gragnic@capgemini.com"),
]
for i, (nom, role, mail) in enumerate(membres):
    x = Inches(0.7) + i * Inches(4.2)
    add_rect(s8, x, Inches(2.5), Inches(3.8), Inches(2.2), RGBColor(0x2D, 0x3B, 0x55))
    add_text(s8, nom,  x + Inches(0.2), Inches(2.65), Inches(3.4), Inches(0.45),
             size=15, bold=True, color=WHITE)
    add_text(s8, role, x + Inches(0.2), Inches(3.1),  Inches(3.4), Inches(0.35),
             size=12, color=RGBColor(0xBC, 0xD9, 0xED))
    add_text(s8, mail, x + Inches(0.2), Inches(3.5),  Inches(3.4), Inches(0.32),
             size=10, color=GRAY)

add_text(s8,
         "Référent AI Office & GTM : Philippe Larroye",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.35),
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s8, "Merci.",
         Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.9),
         size=48, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

add_text(s8,
         "Bureau Veritas · Apave · RATP  —  ISO 9001:2015  —  Claude Opus via GEP Capgemini",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35),
         size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ── Écriture du fichier ────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__),
                        "..", "docs", "UC28_Pitch_Deck_Opus48.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"Fichier généré : {out_path}")
